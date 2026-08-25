import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)        # Slate 900 #0f172a
    CARD_BG = RGBColor(30, 41, 59)        # Slate 800 #1e293b
    ACCENT_RED = RGBColor(255, 71, 87)    # Crimson Red #ff4757
    ACCENT_AMBER = RGBColor(245, 158, 11)  # Amber #f59e0b
    ACCENT_BLUE = RGBColor(99, 102, 241)   # Indigo/Blue #6366f1
    ACCENT_GREEN = RGBColor(46, 204, 113)  # Emerald Green #2ecc71
    TEXT_WHITE = RGBColor(255, 255, 255)  # White
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="에이닷(A.dot) 영어학원 성장의 3대 필수 조건"):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(12)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_RED
        p_cat.font.name = "Malgun Gothic"

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(25)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = "Malgun Gothic"

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.933), Inches(5.1))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = ACCENT_RED
    card1.line.width = Pt(2)

    tf1 = card1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.8)
    tf1.margin_top = Inches(0.9)

    p1 = tf1.paragraphs[0]
    p1.text = "🎯 학원 성장의 3대 필수 조건과 데이터 상권 분석"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_RED
    p1.font.name = "Malgun Gothic"

    p2 = tf1.add_paragraph()
    p2.text = "콘텐츠 + 우수 강사진 + 스마트 상권 분석 웹앱 (서울대 TOP30 & 4년제 대학)"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.font.name = "Malgun Gothic"
    p2.space_before = Pt(20)

    p3 = tf1.add_paragraph()
    p3.text = "에이닷 영어학원 지속 성장 & 신규 출점 시뮬레이션 시스템 보고서 | 2026. 08"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Malgun Gothic"
    p3.space_before = Pt(45)

    # ==========================================
    # SLIDE 2: 학원 성장의 3대 필수 조건
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "1. 학원 성장의 3대 필수 조건 (3 Pillars of Growth)")

    pillars = [
        ("📚 1. 콘텐츠", "• 커리큘럼 및 교재\n• 자체 학습 시스템\n• 1:1 맞춤형 코칭 프로그램\n\n👉 에이닷 독자 브랜드 가치", ACCENT_AMBER),
        ("🎓 2. 우수 강사진", "• 학생 지도 및 강의력\n• 동기부여 및 학습 관리\n• 4년제 대학 인재풀 연동\n\n👉 성과 창출의 핵심 동력", ACCENT_BLUE),
        ("🗺️ 3. 상권 분석", "• [핵심! 이 웹앱이 필요한 이유]\n• 수강생/잠재고객/학원수 정밀 분석\n• 서울대 TOP30 및 학구열 판정\n\n👉 실패 없는 출점의 완성", ACCENT_RED)
    ]

    for idx, (p_title, p_desc, p_color) in enumerate(pillars):
        left_pos = Inches(0.8 + idx * 3.95)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = p_color
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.4)

        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = p_color
        p.font.name = "Malgun Gothic"

        p2 = tf.add_paragraph()
        p2.text = p_desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_WHITE
        p2.font.name = "Malgun Gothic"
        p2.space_before = Pt(16)

    # ==========================================
    # SLIDE 3: 상권 분석의 4대 핵심 데이터 요소
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "2. 상권 분석을 완벽하게 수행하는 4대 데이터 요소")

    data_elements = [
        ("🏫 1. 학생수 추이", "2024~2026년 3개년 학교별 총원 및 학년별 유입/유출 추이 정밀 시각화", ACCENT_RED),
        ("🎯 2. 잠재 고객수", "반경 3km 중·고등학생 총원의 5%를 잠정 고객수로 자동 계산 (유효 타겟)", ACCENT_AMBER),
        ("🏢 3. 아파트 세대수", "반경 3km 내 대단지 주거 타운 세대수를 집계하여 탄탄한 배후 수요 측정", ACCENT_GREEN),
        ("📚 4. 학원수 (밀집도)", "지번별 경쟁 학원수 분포 분석으로 독점적 초희소형(50개 미만) 시장 발굴", ACCENT_BLUE)
    ]

    for idx, (d_title, d_desc, d_color) in enumerate(data_elements):
        row = idx // 2
        col = idx % 2
        left_pos = Inches(0.8 + col * 5.95)
        top_pos = Inches(1.8 + row * 2.5)

        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.6), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = d_color
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = d_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = d_color
        p.font.name = "Malgun Gothic"

        p2 = tf.add_paragraph()
        p2.text = d_desc
        p2.font.size = Pt(13.5)
        p2.font.color.rgb = TEXT_WHITE
        p2.font.name = "Malgun Gothic"
        p2.space_before = Pt(10)

    # ==========================================
    # SLIDE 4: 서울대 입결 TOP 30 - 직관적 학구열 지표
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "3. ⭐ 서울대 입결 TOP 30 고등학교 - 직관적 학구열 파악")

    card_l = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = ACCENT_AMBER
    card_l.line.width = Pt(2)

    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.4)
    tf_l.margin_top = Inches(0.4)

    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "⭐ 서울대 TOP 30 시각화의 의의"
    p_l1.font.size = Pt(20)
    p_l1.font.bold = True
    p_l1.font.color.rgb = ACCENT_AMBER
    p_l1.font.name = "Malgun Gothic"

    p_l2 = tf_l.add_paragraph()
    p_l2.text = (
        "• 고등학교 중 '서울대 합격자수 상위 30개교' 마커 표시\n\n"
        "• 지역 학구열의 직관적 파악:\n"
        "  - 지도에서 해당 지역에 TOP 30 고등학교가 위치해 있는지 한눈에 즉시 확인.\n\n"
        "  - 학구열이 높은 지역일수록 학부모의 교육 투자 성향이 강하고 고난도 영어 코칭(에이닷) 수요가 폭발적임.\n\n"
        "  - 프리미엄 입지 판정의 단초 역할 수행."
    )
    p_l2.font.size = Pt(13.5)
    p_l2.font.color.rgb = TEXT_WHITE
    p_l2.font.name = "Malgun Gothic"
    p_l2.space_before = Pt(14)

    card_r = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = ACCENT_RED
    card_r.line.width = Pt(2)

    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.4)
    tf_r.margin_top = Inches(0.4)

    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "📈 학구열과 에이닷 매출의 상관관계"
    p_r1.font.size = Pt(20)
    p_r1.font.bold = True
    p_r1.font.color.rgb = ACCENT_RED
    p_r1.font.name = "Malgun Gothic"

    p_r2 = tf_r.add_paragraph()
    p_r2.text = (
        "1️⃣ 고학구열 지역의 수강생 LTV (평생가치) 증대:\n"
        "    내신 및 수능 대세를 이끄는 명문고 배후지로 장기 수강률 확보.\n\n"
        "2️⃣ 마케팅 효율 및 입소문 확산 가속:\n"
        "    학부모 커뮤니티가 활성화되어 입소문을 통한 성과 전파가 매우 빠름.\n\n"
        "3️⃣ 에이닷 브랜드 고급화 시너지:\n"
        "    명문고 상권 입점으로 지역 내 대표 영어학원으로 포지셔닝."
    )
    p_r2.font.size = Pt(13.5)
    p_r2.font.color.rgb = TEXT_WHITE
    p_r2.font.name = "Malgun Gothic"
    p_r2.space_before = Pt(14)

    # ==========================================
    # SLIDE 5: 🏛️ 4년제 대학교 - 학구열 + 강사 수급의 2대 가치
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "4. 🏛️ 4년제 대학교의 2대 가치: 학구열 판단 + 강사 수급 용이성")

    card_u1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_u1.fill.solid()
    card_u1.fill.fore_color.rgb = CARD_BG
    card_u1.line.color.rgb = ACCENT_BLUE
    card_u1.line.width = Pt(2)

    tf_u1 = card_u1.text_frame
    tf_u1.word_wrap = True
    tf_u1.margin_left = Inches(0.4)
    tf_u1.margin_top = Inches(0.4)

    p_u1_1 = tf_u1.paragraphs[0]
    p_u1_1.text = "🎓 1. 지역 학구열 지표로서의 가치"
    p_u1_1.font.size = Pt(20)
    p_u1_1.font.bold = True
    p_u1_1.font.color.rgb = ACCENT_BLUE
    p_u1_1.font.name = "Malgun Gothic"

    p_u1_2 = tf_u1.add_paragraph()
    p_u1_2.text = (
        "• 4년제 대학교 입지 지역의 문화적 특징:\n"
        "  - 지적 분위기와 고학력 인구 비중이 높아 지역 전체의 학구열 수준 파악 가능.\n\n"
        "  - 학부모들의 대학 입시 관심도가 높아 고급 인강/학원 코칭 수요가 지속적으로 높음."
    )
    p_u1_2.font.size = Pt(14)
    p_u1_2.font.color.rgb = TEXT_WHITE
    p_u1_2.font.name = "Malgun Gothic"
    p_u1_2.space_before = Pt(16)

    card_u2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_u2.fill.solid()
    card_u2.fill.fore_color.rgb = CARD_BG
    card_u2.line.color.rgb = ACCENT_GREEN
    card_u2.line.width = Pt(2)

    tf_u2 = card_u2.text_frame
    tf_u2.word_wrap = True
    tf_u2.margin_left = Inches(0.4)
    tf_u2.margin_top = Inches(0.4)

    p_u2_1 = tf_u2.paragraphs[0]
    p_u2_1.text = "🤝 2. 강사 구인 용이성 (수급 원활도)"
    p_u2_1.font.size = Pt(20)
    p_u2_1.font.bold = True
    p_u2_1.font.color.rgb = ACCENT_GREEN
    p_u2_1.font.name = "Malgun Gothic"

    p_u2_2 = tf_u2.add_paragraph()
    p_u2_2.text = (
        "• 우수한 강사 수급 리스크 원천 차단:\n"
        "  - 명문 대학생 및 대학원생 파트타임/전임 강사를 가까운 인근에서 수월하게 구인.\n\n"
        "  - 강사 수급 난항으로 인한 개원 지연 및 휴원 위험을 방지하여 에이닷 네트워크 성장의 병목 해결."
    )
    p_u2_2.font.size = Pt(14)
    p_u2_2.font.color.rgb = TEXT_WHITE
    p_u2_2.font.name = "Malgun Gothic"
    p_u2_2.space_before = Pt(16)

    # ==========================================
    # SLIDE 6: 🎯 RDB 추천입지 모델 & 경영진 최종 결론
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "5. 🎯 RDB 추천입지 모델 & 경영진 최종 결론")

    card_fin = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
    card_fin.fill.solid()
    card_fin.fill.fore_color.rgb = CARD_BG
    card_fin.line.color.rgb = ACCENT_RED
    card_fin.line.width = Pt(2)

    tf_fin = card_fin.text_frame
    tf_fin.word_wrap = True
    tf_fin.margin_left = Inches(0.5)
    tf_fin.margin_top = Inches(0.4)

    p_f1 = tf_fin.paragraphs[0]
    p_f1.text = "💡 본 웹앱은 에이닷 성장의 3대 축 중 '상권 분석'의 최첨단 표준 시뮬레이터입니다."
    p_f1.font.size = Pt(20)
    p_f1.font.bold = True
    p_f1.font.color.rgb = ACCENT_RED
    p_f1.font.name = "Malgun Gothic"

    p_f2 = tf_fin.add_paragraph()
    p_f2.text = (
        "• 콘텐츠와 우수 강사진이 준비되어 있다면, 마지막 성공 열쇠는 '상권 분석'입니다.\n\n"
        "• 본 웹앱은 수강생/잠재고객/세대수/학원수 분석은 물론, ⭐ 서울대 TOP30과 🏛️ 4년제 대학교 데이터로 '학구열'과 '강사 구인 용이성'까지 직관적으로 판정합니다.\n\n"
        "• 기존 지점 3km 이격으로 자기잠식을 막고 독점적 🔥 초희소형(학원<50개) 지역을 우선 진출하여 실패 없는 전국 1위 성장을 견인할 것을 제언합니다."
    )
    p_f2.font.size = Pt(14.5)
    p_f2.font.color.rgb = TEXT_WHITE
    p_f2.font.name = "Malgun Gothic"
    p_f2.space_before = Pt(16)

    # Save presentations with fallback if open in PowerPoint
    out1_a = r"C:\Users\PC-B-088\.gemini\antigravity\scratch\adotinformationmap\adot_expansion_presentation.pptx"
    out1_b = r"C:\Users\PC-B-088\.gemini\antigravity\scratch\adotinformationmap\adot_expansion_presentation_v2.pptx"
    out2_a = r"C:\Users\PC-B-088\.gemini\antigravity\brain\0af80624-cd74-465c-8de5-45860e05e8b4\adot_expansion_presentation.pptx"
    out2_b = r"C:\Users\PC-B-088\.gemini\antigravity\brain\0af80624-cd74-465c-8de5-45860e05e8b4\adot_expansion_presentation_v2.pptx"
    
    try:
        prs.save(out1_a)
        target1 = out1_a
    except Exception:
        prs.save(out1_b)
        target1 = out1_b

    try:
        prs.save(out2_a)
        target2 = out2_a
    except Exception:
        prs.save(out2_b)
        target2 = out2_b

    print("PowerPoint presentation created successfully:")
    print("   - Project path:", target1)
    print("   - Artifact path:", target2)

if __name__ == "__main__":
    create_presentation()
